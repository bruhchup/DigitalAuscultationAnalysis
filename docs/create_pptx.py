"""Generate AusculTek progress report PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Theme colours --
DARK_TEAL = RGBColor(0x0C, 0x4A, 0x6E)
MID_TEAL = RGBColor(0x08, 0x91, 0xB2)
LIGHT_BG = RGBColor(0xF0, 0xF9, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── helpers ──────────────────────────────────────────────────────────────

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=DARK_TEAL, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_bullet_slide_content(tf, items, size=16, color=DARK_TEAL):
    """Add bulleted items to an existing text frame."""
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0


def title_bar(slide, title_text, subtitle_text=None):
    """Dark teal bar across top with title."""
    add_shape(slide, 0, 0, W, Inches(1.4), fill=DARK_TEAL)
    t = tb(slide, Inches(0.6), Inches(0.25), Inches(11), Inches(0.7))
    set_text(t.text_frame, title_text, size=32, color=WHITE, bold=True)
    if subtitle_text:
        t2 = tb(slide, Inches(0.6), Inches(0.85), Inches(11), Inches(0.4))
        set_text(t2.text_frame, subtitle_text, size=16, color=ACCENT)
    # accent line
    add_shape(slide, 0, Inches(1.4), W, Inches(0.06), fill=MID_TEAL)


def footer(slide):
    t = tb(slide, Inches(0.4), Inches(7.0), Inches(5), Inches(0.4))
    set_text(t.text_frame, "AusculTek  |  CIS 487 Capstone  |  Spring 2026", size=10, color=GRAY)


# ── SLIDE 1 : Title ─────────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, DARK_TEAL)

# big title
t = tb(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.2))
set_text(t.text_frame, "AusculTek", size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# subtitle
t = tb(sl, Inches(1), Inches(3.0), Inches(11), Inches(0.8))
set_text(t.text_frame, "Automated Respiratory Sound Classification Using Deep Learning",
         size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

# accent line
add_shape(sl, Inches(4.5), Inches(3.9), Inches(4), Inches(0.04), fill=MID_TEAL)

# author info
t = tb(sl, Inches(1), Inches(4.3), Inches(11), Inches(1.2))
tf = t.text_frame
tf.word_wrap = True
for line in ["Hayden Banks", "CIS 487 -- Capstone Progress Report", "April 3, 2026"]:
    if tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xBA, 0xE6, 0xFD)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)


# ── SLIDE 2 : Problem & Motivation ──────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Problem & Motivation")

items = [
    "Respiratory diseases cause 3M+ deaths/year (COPD alone) -- early detection is critical",
    "Auscultation (stethoscope listening) is subjective -- inter-listener agreement < 75%",
    "No permanent record -- sounds can't be revisited, shared, or tracked over time",
    "Digital stethoscopes create .wav files -- enabling computational analysis",
    "Goal: Build an AI assistant that provides objective, repeatable lung sound classification",
]
t = tb(sl, Inches(0.6), Inches(1.8), Inches(12), Inches(5))
add_bullet_slide_content(t.text_frame, items, size=20)
footer(sl)


# ── SLIDE 3 : Project Overview ───────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Project Overview", "End-to-end respiratory sound classification system")

# 4 pipeline boxes
labels = ["Audio\nPreprocessing", "Segmentation", "Classification", "Visualization"]
descs = [
    "Resample to 16 kHz\nBandpass filter 50-2000 Hz",
    "Sliding window or\nICBHI annotations",
    "CNN6 + SCL\nor Random Forest",
    "Streamlit dashboard\nInteractive charts",
]
colors = [MID_TEAL, RGBColor(0x06, 0xB6, 0xD4), RGBColor(0x0E, 0xA5, 0xE9), DARK_TEAL]

box_w = Inches(2.7)
gap = Inches(0.3)
start_x = Inches(0.6)
for i in range(4):
    x = start_x + i * (box_w + gap)
    # box
    s = add_shape(sl, x, Inches(2.2), box_w, Inches(2.6), fill=colors[i])
    s.shadow.inherit = False
    # label
    t = tb(sl, x, Inches(2.4), box_w, Inches(0.9))
    set_text(t.text_frame, labels[i], size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # desc
    t = tb(sl, x, Inches(3.3), box_w, Inches(1.2))
    set_text(t.text_frame, descs[i], size=14, color=RGBColor(0xE0, 0xF2, 0xFE), alignment=PP_ALIGN.CENTER)

    # arrow between boxes
    if i < 3:
        ax = x + box_w + Inches(0.02)
        t = tb(sl, ax, Inches(3.1), Inches(0.3), Inches(0.5))
        set_text(t.text_frame, ">", size=28, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

footer(sl)


# ── SLIDE 4 : Dataset ────────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Dataset: ICBHI 2017", "International Conference on Biomedical and Health Informatics")

left_items = [
    "920 recordings from 126 patients",
    "5.5 hours of respiratory audio",
    "Expert-annotated respiratory cycle boundaries",
    "4 classes: Normal, Crackle, Wheeze, Both",
    "Standardized 60/40 train/test split",
]
t = tb(sl, Inches(0.6), Inches(1.8), Inches(6), Inches(4))
add_bullet_slide_content(t.text_frame, left_items, size=18)

# class distribution table
headers = ["Class", "Samples", "Proportion"]
rows = [
    ("Normal", "2,063", "49.8%"),
    ("Crackle", "1,215", "29.3%"),
    ("Wheeze", "501", "12.1%"),
    ("Both", "363", "8.8%"),
]

table_shape = sl.shapes.add_table(5, 3, Inches(7.5), Inches(2.0), Inches(4.5), Inches(2.8))
table = table_shape.table
table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(1.5)

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_TEAL
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEAL
            p.alignment = PP_ALIGN.CENTER

footer(sl)


# ── SLIDE 5 : Model Architecture ─────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Model Architecture", "CNN6 with Supervised Contrastive Learning")

left_items = [
    "CNN6 backbone (PANNs family)",
    "  - 4 convolutional blocks, 5x5 kernels",
    "  - Channels: 64 > 128 > 256 > 512",
    "  - Pretrained on AudioSet (2M clips, 527 classes)",
    "",
    "Hybrid SCL + CE training objective",
    "  - L = 0.5 * L_SupCon + 0.5 * L_CE",
    "  - Contrastive loss pulls same-class embeddings together",
    "  - CE loss with inverse-frequency class weights",
    "",
    "Projector head: 512 > 512 > 128 (for contrastive loss)",
    "Linear classifier: 512 > 4 (for predictions)",
]
t = tb(sl, Inches(0.6), Inches(1.8), Inches(7), Inches(5))
add_bullet_slide_content(t.text_frame, left_items, size=16)

# hyperparams box
s = add_shape(sl, Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.5), fill=WHITE)
t = tb(sl, Inches(8.4), Inches(1.9), Inches(4.2), Inches(0.5))
set_text(t.text_frame, "Training Config", size=18, color=DARK_TEAL, bold=True)

params = [
    "Optimizer: Adam (lr=1e-4)",
    "Scheduler: Cosine Annealing",
    "Batch size: 128",
    "Duration: 8s @ 16 kHz",
    "Mel bands: 64",
    "FFT: 1024, Hop: 512",
    "Freq range: 50-2000 Hz",
    "Augmentation: SpecAugment",
    "Temperature (tau): 0.06",
    "GPU: Google Colab T4",
]
t = tb(sl, Inches(8.4), Inches(2.5), Inches(4.2), Inches(3.5))
add_bullet_slide_content(t.text_frame, params, size=14, color=GRAY)
footer(sl)


# ── SLIDE 6 : Training Results ───────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Current Results", "50-epoch hybrid SCL training run")

# metric cards
metrics = [
    ("ICBHI Score", "0.5497", "(SE + SP) / 2"),
    ("Sensitivity", "0.3330", "Abnormal detection rate"),
    ("Specificity", "0.7663", "Normal detection rate"),
    ("Epochs", "50 / 400", "Target: 400 epochs"),
]
card_w = Inches(2.8)
gap = Inches(0.3)
start_x = Inches(0.6)
for i, (label, value, desc) in enumerate(metrics):
    x = start_x + i * (card_w + gap)
    add_shape(sl, x, Inches(1.8), card_w, Inches(2.2), fill=WHITE)
    # label
    t = tb(sl, x, Inches(1.9), card_w, Inches(0.4))
    set_text(t.text_frame, label, size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
    # value
    t = tb(sl, x, Inches(2.3), card_w, Inches(0.7))
    set_text(t.text_frame, value, size=36, color=DARK_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    # desc
    t = tb(sl, x, Inches(3.2), card_w, Inches(0.4))
    set_text(t.text_frame, desc, size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

# analysis
analysis = [
    "Model achieves competitive specificity (76.6%) -- good at identifying normal sounds",
    "Sensitivity needs improvement (33.3%) -- misses some abnormal sounds",
    "Class imbalance (Normal ~50%) biases predictions toward the majority class",
    "Extended training to 400 epochs expected to improve sensitivity significantly",
    "Current score is within range of published deep learning results on ICBHI",
]
t = tb(sl, Inches(0.6), Inches(4.3), Inches(12), Inches(2.5))
add_bullet_slide_content(t.text_frame, analysis, size=16)
footer(sl)


# ── SLIDE 7 : Web Application ────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Web Application", "Interactive Streamlit dashboard")

features = [
    "Upload .wav lung recordings via sidebar",
    "Optional ICBHI annotation file support",
    "Real-time classification with confidence scores",
    "",
    "Dashboard visualizations:",
    "  - Overall confidence gauge chart",
    "  - Classification distribution pie chart",
    "  - Color-coded waveform overlay",
    "  - Class probability heatmap across segments",
    "  - Per-segment confidence bar chart",
    "  - Detailed segment data table",
    "",
    "Session history tracking",
    "Raw JSON output for debugging",
    "Dual-model support (CNN6 or Random Forest)",
]
t = tb(sl, Inches(0.6), Inches(1.8), Inches(6.5), Inches(5))
add_bullet_slide_content(t.text_frame, features, size=16)

# tech stack box
s = add_shape(sl, Inches(7.8), Inches(1.8), Inches(5), Inches(4.2), fill=WHITE)
t = tb(sl, Inches(8.0), Inches(1.9), Inches(4.6), Inches(0.5))
set_text(t.text_frame, "Tech Stack", size=18, color=DARK_TEAL, bold=True)

stack = [
    "Python 3.12",
    "PyTorch 2.0+ / torchaudio",
    "scikit-learn 1.6.1",
    "librosa 0.11.0",
    "Streamlit 1.30+",
    "Plotly 5.18+",
    "Google Colab (T4 GPU)",
    "GitHub for version control",
]
t = tb(sl, Inches(8.0), Inches(2.5), Inches(4.6), Inches(3.2))
add_bullet_slide_content(t.text_frame, stack, size=14, color=GRAY)
footer(sl)


# ── SLIDE 8 : Literature ─────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Key References")

refs = [
    ("Kong et al. (2020)", "PANNs: Large-Scale Pretrained Audio Neural Networks",
     "Provides the CNN6 backbone pretrained on AudioSet -- enables transfer learning to respiratory sounds"),
    ("Moummad & Farrugia (2023)", "Lung Sound Classification Using SCL",
     "Primary reference -- demonstrates SCL outperforms cross-entropy on ICBHI benchmark"),
    ("Gunel et al. (2021)", "Supervised Contrastive Learning for Fine-tuning",
     "Motivated the hybrid SCL+CE training objective used in this project"),
    ("Rocha et al. (2019)", "ICBHI Open Access Database",
     "Defines the dataset, evaluation protocol, and ICBHI score metric"),
]

for i, (author, title, note) in enumerate(refs):
    y = Inches(1.8) + i * Inches(1.3)
    add_shape(sl, Inches(0.6), y, Inches(12), Inches(1.15), fill=WHITE)
    t = tb(sl, Inches(0.8), y + Inches(0.08), Inches(11.5), Inches(0.4))
    set_text(t.text_frame, f"{author}  --  {title}", size=16, color=DARK_TEAL, bold=True)
    t = tb(sl, Inches(0.8), y + Inches(0.55), Inches(11.5), Inches(0.5))
    set_text(t.text_frame, note, size=14, color=GRAY)

footer(sl)


# ── SLIDE 9 : Future Work ────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
title_bar(sl, "Future Work")

future = [
    ("Extended Training", "Complete 400-epoch run on Colab T4 -- notebook is ready to go"),
    ("Hyperparameter Tuning", "Adjust class weights, tau, alpha to improve sensitivity"),
    ("Backbone Comparison", "Evaluate CNN10 and CNN14 -- code already supports all three"),
    ("Real-Time Inference", "Stream from connected digital stethoscope for live analysis"),
    ("Clinical Validation", "Partner with healthcare professionals for real-world evaluation"),
    ("Model Explainability", "Add Grad-CAM visualizations to show what drives predictions"),
]

for i, (title, desc) in enumerate(future):
    y = Inches(1.8) + i * Inches(0.85)
    # accent dot
    add_shape(sl, Inches(0.7), y + Inches(0.12), Inches(0.15), Inches(0.15), fill=MID_TEAL)
    t = tb(sl, Inches(1.1), y, Inches(3.5), Inches(0.4))
    set_text(t.text_frame, title, size=18, color=DARK_TEAL, bold=True)
    t = tb(sl, Inches(4.6), y, Inches(8), Inches(0.4))
    set_text(t.text_frame, desc, size=16, color=GRAY)

footer(sl)


# ── SLIDE 10 : Thank You ─────────────────────────────────────────────────

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, DARK_TEAL)

t = tb(sl, Inches(1), Inches(2.2), Inches(11), Inches(1))
set_text(t.text_frame, "Thank You", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(sl, Inches(5), Inches(3.4), Inches(3), Inches(0.04), fill=MID_TEAL)

t = tb(sl, Inches(1), Inches(3.8), Inches(11), Inches(1.5))
tf = t.text_frame
tf.word_wrap = True
for line in ["Hayden Banks", "github.com/bruhchup/DigitalAuscultationAnalysis", "Questions?"]:
    if tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xBA, 0xE6, 0xFD)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)


# ── Save ──────────────────────────────────────────────────────────────────

out = r"E:\Shepherd\Spring2025\CIS487\personal_branch\Sentiment-Analysis-WLB\DigitalAuscultationAnalysis\docs\AusculTek_Progress_Report.pptx"
prs.save(out)
print(f"Saved: {out}")
